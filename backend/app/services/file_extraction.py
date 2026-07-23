import io

from fastapi import HTTPException

SUPPORTED_EXTENSIONS = (".pdf", ".pptx")


def extract_text(filename: str, file_bytes: bytes) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    if lower.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    raise HTTPException(
        status_code=400,
        detail=f"Unsupported file type. Supported: {', '.join(SUPPORTED_EXTENSIONS)}",
    )


def _extract_pdf(file_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise HTTPException(status_code=500, detail="PDF support isn't installed on the server.")

    reader = PdfReader(io.BytesIO(file_bytes))
    pages_text = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages_text if p.strip())
    if not text:
        raise HTTPException(
            status_code=400,
            detail="No extractable text found in this PDF (it may be scanned/image-only).",
        )
    return text


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=500, detail="PowerPoint support isn't installed on the server.")

    presentation = Presentation(io.BytesIO(file_bytes))
    slide_texts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            slide_texts.append(f"Slide {slide_index}:\n" + "\n".join(parts))

    text = "\n\n".join(slide_texts)
    if not text:
        raise HTTPException(status_code=400, detail="No extractable text found in this presentation.")
    return text
